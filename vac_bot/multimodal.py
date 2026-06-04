import os
import base64
import csv
from pathlib import Path
from openai import OpenAI

client = OpenAI()


def extract_image_text(image_path: str) -> str:
    """Use GPT-4o vision to describe an image in detail."""
    with open(image_path, "rb") as f:
        b64 = base64.b64encode(f.read()).decode("utf-8")

    ext = Path(image_path).suffix.lower().lstrip(".")
    if ext == "jpg":
        ext = "jpeg"
    mime = f"image/{ext}"

    resp = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": (
                            "Describe this image in detail. Extract all visible text, "
                            "identify any diagrams, charts, graphs, screenshots, or UI elements. "
                            "Be thorough — every piece of information matters for a knowledge base."
                        ),
                    },
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:{mime};base64,{b64}", "detail": "high"},
                    },
                ],
            }
        ],
        max_tokens=4096,
    )
    return resp.choices[0].message.content or ""


def extract_table_text(file_path: str) -> str:
    """Parse xlsx/csv and return markdown-formatted tables."""
    ext = Path(file_path).suffix.lower()
    parts = []

    if ext == ".csv":
        with open(file_path, "r", encoding="utf-8-sig") as f:
            reader = csv.reader(f)
            rows = list(reader)
        if rows:
            parts.append(_rows_to_markdown(rows))

    elif ext == ".xlsx":
        import openpyxl

        wb = openpyxl.load_workbook(file_path, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cleaned = [str(c) if c is not None else "" for c in row]
                if any(v.strip() for v in cleaned):
                    rows.append(cleaned)
            if rows:
                parts.append(f"## Sheet: {sheet_name}\n\n")
                parts.append(_rows_to_markdown(rows))
                parts.append("\n\n")

    return "\n".join(parts).strip()


def _rows_to_markdown(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    header = "| " + " | ".join(str(c) for c in rows[0]) + " |"
    sep = "| " + " | ".join("---" for _ in rows[0]) + " |"
    body_lines = [
        "| " + " | ".join(str(c) for c in row) + " |" for row in rows[1:]
    ]
    body = "\n".join(body_lines) if body_lines else ""
    return header + "\n" + sep + ("\n" + body if body else "")


def extract_slides_text(file_path: str) -> list[dict]:
    """Extract text per slide from a .pptx file."""
    from pptx import Presentation

    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    texts.append(" | ".join(cell.text.strip() for cell in row.cells))
        if texts:
            slides.append({"slide_number": i, "text": "\n".join(texts)})
    return slides


def extract_scanned_pdf_text(file_path: str) -> list[dict]:
    """OCR scanned/handwritten PDF pages using GPT-4o vision."""
    import fitz

    doc = fitz.open(file_path)
    pages = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        pix = page.get_pixmap(dpi=300)
        img_bytes = pix.tobytes("png")
        b64 = base64.b64encode(img_bytes).decode("utf-8")

        resp = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": (
                                "Extract all text from this document page exactly as written. "
                                "Preserve the original structure, paragraphs, lists, and formatting. "
                                "If there are tables, output them as markdown tables. "
                                "If there are handwritten notes, transcribe them as accurately as possible."
                            ),
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{b64}",
                                "detail": "high",
                            },
                        },
                    ],
                }
            ],
            max_tokens=4096,
        )
        text = resp.choices[0].message.content or ""
        if text.strip():
            pages.append({"page_number": page_num + 1, "text": text})
    doc.close()
    return pages
